import argparse
import json
import os
import re
import shutil
import sys
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import Iterable, List, Optional, Sequence, Tuple


@dataclass(frozen=True)
class Rule:
    name: str
    pattern: str
    severity: str
    extensions: Tuple[str, ...]


@dataclass(frozen=True)
class Finding:
    project: str
    dir: str
    file: str
    rule: str
    severity: str
    hits: int
    lines: str


def normalize_severity(v: object) -> str:
    s = str(v or "").strip().lower()
    if s in {"high", "h", "p0", "0", "严重", "高"}:
        return "high"
    if s in {"medium", "m", "p1", "1", "中"}:
        return "medium"
    if s in {"low", "l", "p2", "2", "低"}:
        return "low"
    return "medium"


def parse_extensions(v: object) -> Tuple[str, ...]:
    raw = str(v or "").strip()
    if not raw:
        return tuple()
    parts = re.split(r"[,，;\s]+", raw)
    out: List[str] = []
    ignore_words = {"提交件", "工作件", ".提交件", ".工作件"}
    for p in parts:
        p = p.strip()
        if not p:
            continue
        if p in ignore_words:
            continue
        if not p.startswith("."):
            p = "." + p
        p = p.lower()
        if p in ignore_words:
            continue
        if not re.fullmatch(r"\.[a-z0-9]+", p):
            continue
        out.append(p)
    return tuple(out)


def load_rules_from_excel(path: Path) -> List[Rule]:
    try:
        import openpyxl  # type: ignore
    except Exception as e:  # noqa: BLE001
        raise RuntimeError("Missing dependency: openpyxl. Install via: pip install openpyxl") from e

    wb = openpyxl.load_workbook(path, data_only=True)

    def norm_key(s: object) -> str:
        return re.sub(r"[\s_\-()（）【】\[\].,:，;；/\\]+", "", str(s or "").strip().lower())

    def guess_kind(header: str) -> str:
        k = norm_key(header)
        if not k:
            return ""
        if any(x in k for x in ("编号", "序号", "id", "no", "index")):
            return "id"
        if any(x in k for x in ("严重", "等级", "level", "severity", "priority", "优先")):
            return "severity"
        if any(x in k for x in ("后缀", "suffix", "ext", "extensions", "文件类型", "filetype")):
            return "extensions"
        if any(x in k for x in ("对象", "scope", "target", "范围", "文件名", "filename")):
            return "scope"
        if any(x in k for x in ("名称", "标题", "name", "title", "规则名")):
            return "name"
        if any(
            x in k
            for x in (
                "规则",
                "正则",
                "regex",
                "regexp",
                "keyword",
                "关键字",
                "表达式",
                "pattern",
                "检查内容",
                "检查项",
                "检测点",
                "要点",
                "校验项",
                "匹配内容",
            )
        ):
            return "pattern"
        return ""

    def is_numeric_like(v: object) -> bool:
        s = str(v or "").strip()
        if not s:
            return False
        if re.fullmatch(r"\d+(\.\d+)?", s):
            return True
        return False

    def is_non_rule_marker(v: object) -> bool:
        s = str(v or "").strip().lower()
        if not s:
            return False
        return s in {
            "-",
            "—",
            "–",
            "不符合",
            "符合",
            "执行总数",
            "检查项",
            "检查结果",
            "不适用",
            "y",
            "yes",
            "true",
            "1",
            "是",
            "对",
            "√",
            "✓",
            "✔",
            "n",
            "no",
            "false",
            "0",
            "否",
            "错",
            "×",
            "✗",
            "✘",
        }

    def score_checkpoint_col(rows: List[Sequence[object]], header_i: int, headers: List[str], kinds: List[str], id_cols: set) -> int:
        for i, h in enumerate(headers):
            if i in id_cols:
                continue
            hk = norm_key(h)
            if not hk:
                continue
            if "检查要点" in hk or "检测要点" in hk or "校验要点" in hk:
                return i
            if "检查点" in hk or "检测点" in hk or "要点" in hk:
                return i

        best_idx = -1
        best_score = -1.0
        max_rows = min(len(rows), header_i + 1 + 30)
        sev_set = {"high", "medium", "low", "h", "m", "l", "p0", "p1", "p2", "严重", "高", "中", "低"}
        scope_set = {"file", "filename", "name", "文件名", "文件", "content", "内容", "正文"}

        for c in range(len(headers)):
            if c in id_cols:
                continue
            h = headers[c]
            k = kinds[c] if c < len(kinds) else ""
            header_score = 0.0
            hk = norm_key(h)
            if k == "pattern":
                header_score += 6.0
            elif any(x in hk for x in ("要点", "检查", "检测", "校验", "匹配", "规则", "内容", "表达式", "正则", "关键字")):
                header_score += 3.0

            n = 0
            numeric = 0
            bool_mark = 0
            sev = 0
            ext = 0
            scope = 0
            total_len = 0
            for r in range(header_i + 1, max_rows):
                row = rows[r]
                if c >= len(row):
                    continue
                v = as_text(row[c])
                if not v:
                    continue
                n += 1
                total_len += len(v)
                if is_numeric_like(v):
                    numeric += 1
                if is_non_rule_marker(v):
                    bool_mark += 1
                vl = v.lower()
                if vl in sev_set:
                    sev += 1
                if vl in scope_set:
                    scope += 1
                if re.fullmatch(r"\.[a-z0-9]+", vl) or re.fullmatch(r"[a-z0-9]+", vl):
                    t = vl if vl.startswith(".") else "." + vl
                    if re.fullmatch(r"\.[a-z0-9]+", t):
                        ext += 1

            if n <= 0:
                continue
            avg_len = total_len / n
            numeric_ratio = numeric / n
            bool_ratio = bool_mark / n
            sev_ratio = sev / n
            ext_ratio = ext / n
            scope_ratio = scope / n

            value_score = 0.0
            if avg_len >= 4 and numeric_ratio < 0.3 and bool_ratio < 0.3 and sev_ratio < 0.3 and ext_ratio < 0.3 and scope_ratio < 0.3:
                value_score += 6.0
            elif avg_len >= 3 and numeric_ratio < 0.4 and bool_ratio < 0.4:
                value_score += 3.0

            penalty = 0.0
            if sev_ratio >= 0.4 or ext_ratio >= 0.4 or scope_ratio >= 0.4:
                penalty += 4.0
            if numeric_ratio >= 0.6:
                penalty += 6.0

            score = header_score + value_score - penalty
            if score > best_score:
                best_score = score
                best_idx = c

        return best_idx

    def find_header_row(rows: List[Sequence[object]]) -> int:
        best_i = -1
        best_score = -1
        best_non_empty = -1
        for i in range(min(30, len(rows))):
            row = rows[i]
            headers = [str(x or "").strip() for x in row]
            kinds = [guess_kind(h) for h in headers if h]
            score = sum(1 for k in kinds if k in {"pattern", "name", "severity", "extensions", "scope"})
            non_empty = sum(1 for h in headers if h)
            if score > best_score or (score == best_score and non_empty > best_non_empty):
                best_i = i
                best_score = score
                best_non_empty = non_empty
        return best_i

    def as_text(v: object) -> str:
        return str(v or "").strip()

    def pick_first_non_empty(row: Sequence[object]) -> str:
        for v in row:
            t = as_text(v)
            if t:
                return t
        return ""

    out: List[Rule] = []
    for ws in wb.worksheets:
        stage = str(ws.title or "").strip()
        all_rows = list(ws.iter_rows(values_only=True))
        if not all_rows:
            continue
        header_idx = find_header_row(all_rows)
        if header_idx < 0:
            header_idx = 0
        header_row = all_rows[header_idx]
        headers = [as_text(x) or f"COL_{i+1}" for i, x in enumerate(header_row)]
        kinds = [guess_kind(h) for h in headers]

        id_cols = {i for i, k in enumerate(kinds) if k == "id"}
        name_col = next((i for i, k in enumerate(kinds) if k == "name"), None)
        severity_col = next((i for i, k in enumerate(kinds) if k == "severity"), None)
        ext_col = next((i for i, k in enumerate(kinds) if k == "extensions"), None)

        checkpoint_col = score_checkpoint_col(all_rows, header_idx, headers, kinds, id_cols)
        if checkpoint_col < 0:
            continue

        for r_idx in range(header_idx + 1, len(all_rows)):
            row = all_rows[r_idx]
            if not row:
                continue
            if checkpoint_col >= len(row):
                continue

            checkpoint = as_text(row[checkpoint_col])
            if not checkpoint or is_numeric_like(checkpoint) or is_non_rule_marker(checkpoint):
                fallback = ""
                for c in range(len(headers)):
                    if c in id_cols:
                        continue
                    if severity_col is not None and c == severity_col:
                        continue
                    if ext_col is not None and c == ext_col:
                        continue
                    if name_col is not None and c == name_col:
                        continue
                    t = as_text(row[c]) if c < len(row) else ""
                    if not t or is_numeric_like(t) or is_non_rule_marker(t):
                        continue
                    if re.fullmatch(r"\.[a-z0-9]+", t.lower()) or re.fullmatch(r"[a-z0-9]+", t.lower()):
                        continue
                    fallback = t
                    break
                if not fallback:
                    continue
                checkpoint = fallback

            base_name = checkpoint
            if name_col is not None and name_col < len(row):
                n = as_text(row[name_col])
                if n and not is_numeric_like(n) and not is_non_rule_marker(n):
                    base_name = n

            severity = normalize_severity(row[severity_col] if severity_col is not None and severity_col < len(row) else None)
            exts = parse_extensions(row[ext_col] if ext_col is not None and ext_col < len(row) else None)
            rule_name = f"{stage} - {base_name}" if stage else base_name
            out.append(Rule(name=rule_name.strip(), pattern=checkpoint.strip(), severity=severity, extensions=exts))

    return out


def file_ext(p: Path) -> str:
    return p.suffix.lower()


def is_text_like(ext: str) -> bool:
    return ext not in {".xlsx", ".xls", ".docx", ".pptx", ".pdf"}


def passes_extensions(p: Path, rule: Rule) -> bool:
    if not rule.extensions:
        return True
    ext = file_ext(p)
    return any(ext == x for x in rule.extensions)


def build_matcher(pattern: str):
    m = re.fullmatch(r"/(.+)/([a-zA-Z]*)", pattern.strip())
    if m:
        expr, flags_raw = m.group(1), m.group(2)
        flags = 0
        if "i" in flags_raw:
            flags |= re.IGNORECASE
        try:
            return ("re", re.compile(expr, flags))
        except re.error:
            return ("text", pattern)
    return ("text", pattern)


def count_matches_text(needle: str, text: str, limit: int = 500) -> int:
    if not needle:
        return 0
    n = 0
    start = 0
    while n < limit:
        idx = text.find(needle, start)
        if idx < 0:
            break
        n += 1
        start = idx + len(needle)
    return n


def count_matches_re(rx: re.Pattern, text: str, limit: int = 500) -> int:
    n = 0
    for _ in rx.finditer(text):
        n += 1
        if n >= limit:
            break
    return n


def line_numbers(text: str, matcher, limit: int = 30) -> List[int]:
    lines = text.splitlines()
    kind, obj = matcher
    out: List[int] = []
    if kind == "re":
        rx: re.Pattern = obj
        flags = rx.flags
        try:
            line_rx = re.compile(rx.pattern, flags & ~re.MULTILINE)
        except re.error:
            return []
        for i, line in enumerate(lines, start=1):
            if len(out) >= limit:
                break
            if line_rx.search(line):
                out.append(i)
        return out

    needle: str = str(obj)
    for i, line in enumerate(lines, start=1):
        if len(out) >= limit:
            break
        if needle and needle in line:
            out.append(i)
    return out


def extract_text(p: Path) -> Optional[str]:
    ext = file_ext(p)
    if is_text_like(ext):
        try:
            return p.read_text(encoding="utf-8", errors="ignore")
        except Exception:  # noqa: BLE001
            return None

    if ext == ".xlsx":
        try:
            import openpyxl  # type: ignore
        except Exception as e:  # noqa: BLE001
            raise RuntimeError("Missing dependency: openpyxl. Install via: pip install openpyxl") from e

        try:
            wb = openpyxl.load_workbook(p, data_only=True)
            parts: List[str] = []
            for ws in wb.worksheets:
                rows: List[str] = []
                for r in ws.iter_rows(values_only=True):
                    row = [str(x).strip() for x in r if x is not None and str(x).strip() != ""]
                    if row:
                        rows.append("\t".join(row))
                if rows:
                    parts.append(f"=== {ws.title} ===\n" + "\n".join(rows))
            return "\n\n".join(parts) if parts else ""
        except Exception:  # noqa: BLE001
            return None

    if ext == ".xls":
        try:
            import xlrd  # type: ignore
        except Exception as e:  # noqa: BLE001
            raise RuntimeError("Missing dependency: xlrd (for .xls). Install via: pip install xlrd") from e
        try:
            book = xlrd.open_workbook(str(p))
            parts: List[str] = []
            for sheet in book.sheets():
                rows: List[str] = []
                for r in range(sheet.nrows):
                    vals = []
                    for c in range(sheet.ncols):
                        v = sheet.cell_value(r, c)
                        s = str(v).strip()
                        if s != "":
                            vals.append(s)
                    if vals:
                        rows.append("\t".join(vals))
                if rows:
                    parts.append(f"=== {sheet.name} ===\n" + "\n".join(rows))
            return "\n\n".join(parts) if parts else ""
        except Exception:  # noqa: BLE001
            return None

    if ext == ".docx":
        try:
            import docx  # type: ignore
        except Exception:  # noqa: BLE001
            docx = None  # type: ignore
        try:
            if docx is not None:
                d = docx.Document(str(p))
                return "\n".join([para.text for para in d.paragraphs if para.text])
        except Exception:  # noqa: BLE001
            pass
        try:
            import zipfile
            import xml.etree.ElementTree as ET

            with zipfile.ZipFile(str(p)) as zf:
                xml_bytes = zf.read("word/document.xml")
            root = ET.fromstring(xml_bytes)
            ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
            paras: List[str] = []
            for para in root.findall(".//w:p", ns):
                texts = [t.text for t in para.findall(".//w:t", ns) if t.text]
                s = "".join(texts).strip()
                if s:
                    paras.append(s)
            return "\n".join(paras)
        except Exception:  # noqa: BLE001
            return None

    if ext == ".pptx":
        try:
            from pptx import Presentation  # type: ignore
        except Exception:  # noqa: BLE001
            Presentation = None  # type: ignore
        try:
            if Presentation is not None:
                prs = Presentation(str(p))
                parts: List[str] = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        text = getattr(shape, "text", None)
                        if text:
                            parts.append(str(text))
                return "\n".join(parts)
        except Exception:  # noqa: BLE001
            pass
        try:
            import zipfile
            import xml.etree.ElementTree as ET

            with zipfile.ZipFile(str(p)) as zf:
                names = [n for n in zf.namelist() if n.startswith("ppt/slides/slide") and n.endswith(".xml")]
                names = sorted(names, key=lambda x: x)
                parts: List[str] = []
                ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                for name in names:
                    root = ET.fromstring(zf.read(name))
                    texts = [t.text for t in root.findall(".//a:t", ns) if t.text]
                    s = "\n".join([x.strip() for x in texts if x.strip()])
                    if s.strip():
                        parts.append(s)
            return "\n\n".join(parts)
        except Exception:  # noqa: BLE001
            return None

    if ext == ".pdf":
        try:
            import pdfplumber  # type: ignore
        except Exception:  # noqa: BLE001
            pdfplumber = None  # type: ignore
        try:
            parts: List[str] = []
            if pdfplumber is not None:
                with pdfplumber.open(str(p)) as pdf:
                    for page in pdf.pages:
                        t = page.extract_text() or ""
                        if t.strip():
                            parts.append(t)
                return "\n\n".join(parts)
        except Exception:  # noqa: BLE001
            parts = []
        try:
            import subprocess

            proc = subprocess.run(
                ["pdftotext", str(p), "-"],
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                check=False,
                text=True,
            )
            txt = (proc.stdout or "").strip()
            if txt:
                return txt
        except Exception:  # noqa: BLE001
            pass
        try:
            try:
                from pypdf import PdfReader  # type: ignore
            except Exception:  # noqa: BLE001
                from PyPDF2 import PdfReader  # type: ignore
            reader = PdfReader(str(p))
            parts = []
            for page in reader.pages:
                t = (page.extract_text() or "").strip()
                if t:
                    parts.append(t)
            return "\n\n".join(parts)
        except Exception:  # noqa: BLE001
            if pdfplumber is None:
                raise RuntimeError(
                    "Missing dependency: pdfplumber or pypdf. Install via: pip install pdfplumber pypdf"
                )
            return None

    return None


def iter_files(root: Path, max_mb: int) -> Iterable[Path]:
    max_bytes = max_mb * 1024 * 1024
    for dirpath, _, filenames in os.walk(root):
        for name in filenames:
            p = Path(dirpath) / name
            try:
                if not p.is_file():
                    continue
                if p.stat().st_size > max_bytes:
                    continue
            except Exception:  # noqa: BLE001
                continue
            yield p


def project_and_relpath(root: Path, p: Path) -> Tuple[str, str, str]:
    rel = p.relative_to(root).as_posix()
    parts = rel.split("/")
    if len(parts) <= 1:
        return ("root", "-", rel)
    proj = parts[0]
    within = "/".join(parts[1:])
    d = "/".join(parts[1:-1]) if len(parts) > 2 else "-"
    return (proj, d or "-", within)


def scan(
    root: Path,
    rules: Sequence[Rule],
    include_exts: Optional[Sequence[str]],
    max_mb: int,
) -> List[Finding]:
    findings: List[Finding] = []
    for p in iter_files(root, max_mb=max_mb):
        ext = file_ext(p)
        if include_exts is not None and ext not in include_exts:
            continue
        if include_exts is None and ext == "":
            continue

        content = extract_text(p)
        if content is None:
            continue

        proj, d, within = project_and_relpath(root, p)
        for r in rules:
            if not passes_extensions(p, r):
                continue
            matcher = build_matcher(r.pattern)
            kind, obj = matcher
            hits = count_matches_re(obj, content) if kind == "re" else count_matches_text(str(obj), content)
            if hits <= 0:
                continue
            lines = ""
            if is_text_like(ext):
                nums = line_numbers(content, matcher)
                lines = ", ".join(str(x) for x in nums)
            findings.append(
                Finding(
                    project=proj,
                    dir=d,
                    file=within,
                    rule=r.name,
                    severity=r.severity,
                    hits=hits,
                    lines=lines,
                )
            )
    return findings


def write_report_xlsx(out_path: Path, findings: Sequence[Finding]) -> None:
    try:
        import openpyxl  # type: ignore
    except Exception as e:  # noqa: BLE001
        raise RuntimeError("Missing dependency: openpyxl. Install via: pip install openpyxl") from e

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "report"
    headers = ["项目", "目录", "文件", "规则", "严重性", "命中次数", "行号"]
    ws.append(headers)

    for x in findings:
        sev = "高" if x.severity == "high" else ("低" if x.severity == "low" else "中")
        ws.append([x.project, x.dir, x.file, x.rule, sev, x.hits, x.lines])

    wb.save(out_path)


def write_rules_json(out_path: Path, rules: Sequence[Rule]) -> None:
    payload = [
        {"name": r.name, "pattern": r.pattern, "severity": r.severity, "extensions": list(r.extensions)} for r in rules
    ]
    out_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def create_api_app(repo_root: Path, config_path: Optional[Path] = None):
    from fastapi import FastAPI, Request
    from fastapi.middleware.cors import CORSMiddleware
    from fastapi.responses import JSONResponse
    from urllib import request as urlrequest
    from urllib.error import HTTPError, URLError
    import logging

    app = FastAPI(title="auto-QA API", version="0.1")
    app.add_middleware(
        CORSMiddleware,
        allow_origins=["*"],
        allow_credentials=True,
        allow_methods=["*"],
        allow_headers=["*"],
    )

    work_dir = (repo_root / "work").resolve()
    work_rules_dir = (work_dir / "rules").resolve()
    work_inputs_dir = (work_dir / "inputs").resolve()
    work_outputs_dir = (work_dir / "outputs").resolve()
    work_scan_dir = (work_dir / "scan").resolve()

    def ensure_work_dirs() -> List[str]:
        paths = [
            str(work_dir),
            str(work_inputs_dir),
            str(work_outputs_dir),
            str(work_rules_dir),
            str(work_scan_dir),
        ]
        work_dir.mkdir(parents=True, exist_ok=True)
        work_inputs_dir.mkdir(parents=True, exist_ok=True)
        work_outputs_dir.mkdir(parents=True, exist_ok=True)
        work_rules_dir.mkdir(parents=True, exist_ok=True)
        work_scan_dir.mkdir(parents=True, exist_ok=True)
        return paths

    ensure_work_dirs()

    skill_prompt_path = (work_inputs_dir / "autoqa.skill.prompt.txt").resolve()

    def load_skill_prompt_template() -> str:
        try:
            return skill_prompt_path.read_text(encoding="utf-8")
        except Exception:  # noqa: BLE001
            return ""

    def norm(s: str) -> str:
        return re.sub(r"[\s_\-()（）【】\[\].,:，;；/\\]+", "", str(s or "").strip().lower())

    def extract_keywords(text: str) -> List[str]:
        t = str(text or "").strip()
        if not t:
            return []
        parts = re.split(r"[\s,，;；/\\|]+", t)
        stop = {"的", "和", "或", "及", "与", "是否", "需要", "必须", "应", "应当", "不得", "禁止", "可以"}
        out: List[str] = []
        for p in parts:
            p = p.strip()
            if len(p) < 2:
                continue
            if p in stop:
                continue
            out.append(p)
        base = norm(t)
        for x in NAMING_SKILL:
            for k in x["keywords"]:
                if norm(k) and norm(k) in base:
                    out.append(k)
        seen = set()
        uniq: List[str] = []
        for x in out:
            k = norm(x)
            if not k or k in seen:
                continue
            seen.add(k)
            uniq.append(x)
        return uniq[:12]

    NAMING_SKILL = [
        {"category": "项目计划", "name": "项目计划-YYYYMMDD-vX.Y.xlsx", "keywords": ["计划", "里程碑", "wbs", "schedule"]},
        {"category": "质量计划", "name": "质量计划-YYYYMMDD-vX.Y.docx / 质量检查清单-YYYYMMDD.xlsx", "keywords": ["质量", "qa", "检查", "审计", "整改"]},
        {"category": "需求文档", "name": "需求规格说明书(PRD)-YYYYMMDD-vX.Y.docx", "keywords": ["需求", "prd", "srs", "规格"]},
        {"category": "设计文档", "name": "概要设计/详细设计-YYYYMMDD-vX.Y.docx", "keywords": ["设计", "hld", "lld", "架构", "接口", "api"]},
        {"category": "测试方案", "name": "测试方案-YYYYMMDD-vX.Y.docx", "keywords": ["测试", "方案", "test", "case"]},
        {"category": "测试报告", "name": "测试报告-YYYYMMDD-vX.Y.docx", "keywords": ["测试", "报告", "result"]},
        {"category": "投产方案", "name": "投产方案-YYYYMMDD-vX.Y.docx", "keywords": ["投产", "上线", "发布", "变更"]},
        {"category": "运维手册", "name": "运维手册-YYYYMMDD-vX.Y.docx", "keywords": ["运维", "手册", "监控", "应急"]},
        {"category": "风险清单", "name": "风险清单-YYYYMMDD.xlsx", "keywords": ["风险", "risk", "问题", "阻塞"]},
        {"category": "会议纪要", "name": "周会纪要-YYYYMMDD.docx / 评审纪要-YYYYMMDD.docx", "keywords": ["会议", "纪要", "评审", "周会"]},
    ]

    def infer_naming_rules(stage: str, checkpoint: str, exts: List[str]) -> List[str]:
        kws = extract_keywords(checkpoint)
        kset = {norm(x) for x in kws}
        base = norm(checkpoint)
        ranked: List[tuple[int, dict]] = []
        for x in NAMING_SKILL:
            score = 0
            for k in x["keywords"]:
                nk = norm(k)
                if nk and (nk in kset or nk in base):
                    score += 1
            if exts:
                name = str(x["name"]).lower()
                ok = any((e.lower() in name) for e in exts if e)
                if not ok:
                    score = 0
            if score > 0:
                ranked.append((score, x))
        ranked.sort(key=lambda t: t[0], reverse=True)
        out: List[str] = []
        for _, x in ranked[:4]:
            out.append(f'{x["category"]}：{x["name"]}（关键词：{"|".join(x["keywords"])}）')
        if out:
            return out
        if exts:
            return [f'未命中命名技能库；请结合检查要点补充更明确的交付物名称线索（后缀：{", ".join(exts)}）']
        return ['未命中命名技能库；请结合检查要点补充更明确的交付物名称线索']

    def infer_naming_rules_meta(stage: str, checkpoint: str, exts: List[str]) -> dict:
        kws = extract_keywords(checkpoint)
        kset = {norm(x) for x in kws}
        base = norm(checkpoint)
        ranked: List[tuple[int, dict]] = []
        for x in NAMING_SKILL:
            score = 0
            for k in x["keywords"]:
                nk = norm(k)
                if nk and (nk in kset or nk in base):
                    score += 1
            if exts:
                name = str(x["name"]).lower()
                ok = any((e.lower() in name) for e in exts if e)
                if not ok:
                    score = 0
            if score > 0:
                ranked.append((score, x))
        ranked.sort(key=lambda t: t[0], reverse=True)
        items: List[dict] = []
        for _, x in ranked[:4]:
            items.append({"category": x["category"], "namePattern": x["name"], "keywords": list(x["keywords"])})
        return {"keywords": kws, "naming": items}

    def generate_prompt_from_skill(rule: dict) -> dict:
        template = load_skill_prompt_template()
        stage = str(rule.get("stage") or "").strip() or "-"
        checkpoint = str(rule.get("checkpoint") or rule.get("pattern") or "").strip()
        scope = "文件名" if str(rule.get("scope") or "").strip() == "file" else "内容"
        sev = str(rule.get("severity") or "medium").strip().lower()
        sev_cn = "高" if sev == "high" else ("低" if sev == "low" else "中")
        exts = rule.get("extensions") or []
        exts = [str(x).strip() for x in exts if str(x).strip()]
        extensions_text = ", ".join(exts) if exts else "不限"

        naming_rules = infer_naming_rules(stage, checkpoint, exts)
        meta = infer_naming_rules_meta(stage, checkpoint, exts)
        kws = extract_keywords(checkpoint)
        target_file_selection = (
            f"- 后缀规则：{extensions_text}\n"
            f"- 命名规则（由技能推断）：\n  " + "\n  ".join(naming_rules) + "\n"
            f"- 关键词规则：从检查要点抽取关键词（如：{'、'.join(kws) or '无'}），要求文件名或章节标题包含关键词\n"
            "- 排除规则：排除统计/结果类表头词（如“符合/不符合/执行总数/检查结果/不适用”等）与明显无关文件\n"
            "- 判定：目标文件集合为空则无法做内容检查，需直接输出缺失交付物"
        )

        check_steps = (
            "- 先执行“目标文件集合推导”得到待检查文件清单\n"
            "- 若目标文件集合为空：输出“不符合：缺少目标交付物文件，无法执行检查”，并给出建议文件名模式与建议目录\n"
            "- 若检查对象=文件名：对每个候选文件验证命名模式/关键词是否满足检查要点，并给出命中文件名作为证据\n"
            "- 若检查对象=内容：在候选文件内容中按检查要点语义检查（必须包含/不得包含/必须说明/结构化要素等），给出命中摘录作为证据"
        )

        output_format = "- 结论：符合 / 不符合\n- 不符合原因：必须可定位（缺文件/缺内容/命名不符）\n- 修改建议：补齐哪个交付物、建议文件名模式、建议目录、建议补充内容位置\n- 证据：文件路径 + 摘录/命中信息"

        if not template.strip():
            prompt = (
                f"【{stage}】\n检查要点（规则详细）：{checkpoint}\n检查对象：{scope}\n严重性：{sev_cn}\n"
                f"适用后缀：{extensions_text}\n\n目标文件集合推导：\n{target_file_selection}\n\n检查步骤：\n{check_steps}\n\n输出格式：\n{output_format}"
            )
        else:
            prompt = template.format(
                stage=stage,
                checkpoint=checkpoint,
                scope=scope,
                severity=sev_cn,
                extensions_text=extensions_text,
                target_file_selection=target_file_selection,
                check_steps=check_steps,
                output_format=output_format,
            )

        return {**rule, "prompt": prompt, "promptMeta": meta}

    config_path = (config_path or (repo_root / "scripts" / "config.yaml")).resolve()

    def parse_yaml_scalar(v: str) -> object:
        s = str(v or "").strip()
        if not s:
            return ""
        if (s.startswith('"') and s.endswith('"')) or (s.startswith("'") and s.endswith("'")):
            return s[1:-1]
        low = s.lower()
        if low in {"true", "yes", "y", "on"}:
            return True
        if low in {"false", "no", "n", "off"}:
            return False
        try:
            if re.fullmatch(r"-?\d+", s):
                return int(s)
            if re.fullmatch(r"-?\d+\.\d+", s):
                return float(s)
        except Exception:  # noqa: BLE001
            return s
        return s

    def load_simple_yaml(p: Path) -> dict:
        if not p.exists() or not p.is_file():
            return {}
        text = p.read_text(encoding="utf-8", errors="replace")
        root: dict = {}
        cur_section: Optional[str] = None
        for raw in text.splitlines():
            if not raw.strip():
                continue
            if raw.lstrip().startswith("#"):
                continue
            indent = len(raw) - len(raw.lstrip(" "))
            line = raw.strip()
            if ":" not in line:
                continue
            if line.endswith(":") and line.count(":") == 1:
                key = line[:-1].strip()
                if not key:
                    continue
                root.setdefault(key, {})
                cur_section = key
                continue
            k, v = line.split(":", 1)
            key = k.strip()
            if not key:
                continue
            val = parse_yaml_scalar(v.strip())
            if indent > 0 and cur_section:
                sec = root.get(cur_section)
                if isinstance(sec, dict):
                    sec[key] = val
                else:
                    root[cur_section] = {key: val}
            else:
                root[key] = val
        return root

    cfg = load_simple_yaml(config_path)
    llm_cfg = cfg.get("llm") if isinstance(cfg, dict) else {}
    llm_cfg = llm_cfg if isinstance(llm_cfg, dict) else {}

    llm_api_key = str(llm_cfg.get("api_key") or "").strip() or str(os.getenv("AUTOQA_LLM_API_KEY") or os.getenv("OPENAI_API_KEY") or "").strip()
    llm_base_url = str(llm_cfg.get("base_url") or "").strip() or str(os.getenv("AUTOQA_LLM_BASE_URL") or os.getenv("OPENAI_BASE_URL") or "https://api.openai.com/v1").strip()
    llm_model = str(llm_cfg.get("model") or "").strip() or str(os.getenv("AUTOQA_LLM_MODEL") or os.getenv("OPENAI_MODEL") or "gpt-4o-mini").strip()
    try:
        llm_timeout = float(llm_cfg.get("timeout") or os.getenv("AUTOQA_LLM_TIMEOUT") or "60")
    except Exception:  # noqa: BLE001
        llm_timeout = 60.0

    llm_base_url = llm_base_url.rstrip("/")
    llm_source = "config.yaml" if str(llm_cfg.get("api_key") or "").strip() else "env"

    def llm_configured() -> bool:
        return bool(str(llm_api_key or "").strip())

    def clamp_text(s: object, max_chars: int) -> str:
        t = str(s or "")
        if max_chars <= 0:
            return ""
        return t[:max_chars]

    def llm_endpoint_url() -> str:
        base = llm_base_url.rstrip("/")
        if base.endswith("/chat/completions"):
            return base
        return base + "/chat/completions"

    def parse_json_maybe(text: str) -> dict:
        t = str(text or "").strip()
        if not t:
            return {}
        try:
            obj = json.loads(t)
            return obj if isinstance(obj, dict) else {}
        except Exception:  # noqa: BLE001
            pass
        i = t.find("{")
        j = t.rfind("}")
        if i >= 0 and j > i:
            try:
                obj = json.loads(t[i : j + 1])
                return obj if isinstance(obj, dict) else {}
            except Exception:  # noqa: BLE001
                return {}
        return {}

    def call_llm(messages: List[dict]) -> str:
        url = llm_endpoint_url()
        payload = {"model": llm_model, "messages": messages, "temperature": 0, "max_tokens": 900}
        data = json.dumps(payload, ensure_ascii=False).encode("utf-8")
        req = urlrequest.Request(url, data=data, method="POST")
        req.add_header("Content-Type", "application/json")
        req.add_header("Authorization", f"Bearer {llm_api_key}")
        try:
            with urlrequest.urlopen(req, timeout=llm_timeout) as resp:
                raw = resp.read().decode("utf-8", errors="replace")
        except HTTPError as e:  # noqa: BLE001
            raise RuntimeError(e.read().decode("utf-8", errors="replace") or str(e)) from e
        except URLError as e:  # noqa: BLE001
            raise RuntimeError(str(e.reason or e)) from e
        obj = json.loads(raw or "{}")
        choices = obj.get("choices") or []
        if not choices:
            raise RuntimeError(raw or "empty llm response")
        msg = (choices[0] or {}).get("message") or {}
        return str(msg.get("content") or "")

    @app.get("/api/health")
    def health():
        return {"ok": True}

    @app.post("/api/work/ensure")
    def ensure_work():
        return {"ok": True, "paths": ensure_work_dirs()}

    @app.get("/api/llm/health")
    def llm_health():
        return {"ok": True, "configured": llm_configured(), "baseUrl": llm_base_url, "model": llm_model, "source": llm_source, "configPath": str(config_path)}

    @app.post("/api/scan/candidates")
    async def scan_candidates(request: Request):
        body = await request.json()
        project = str((body or {}).get("project") or "").strip()
        rule = (body or {}).get("rule") or {}
        files = (body or {}).get("files") or []
        if not isinstance(files, list):
            files = []

        stage_need = str((rule or {}).get("stage") or "").strip()
        checkpoint = str((rule or {}).get("checkpoint") or (rule or {}).get("pattern") or "").strip()
        exts = (rule or {}).get("extensions") or []
        exts = [str(x).strip().lower() for x in exts if str(x).strip()]
        meta = (rule or {}).get("promptMeta") or {}
        naming = meta.get("naming") if isinstance(meta, dict) else []
        naming = naming if isinstance(naming, list) else []
        keywords = meta.get("keywords") if isinstance(meta, dict) else []
        keywords = keywords if isinstance(keywords, list) else []
        naming_patterns = [str((it or {}).get("namePattern") or "").strip() for it in naming if str((it or {}).get("namePattern") or "").strip()]
        hint_text = " ".join([checkpoint] + [str(x or "").strip() for x in (keywords or [])] + naming_patterns).strip()

        def file_stage(x: dict) -> str:
            try:
                return str((x or {}).get("stage") or "").strip()
            except Exception:  # noqa: BLE001
                return ""

        def file_path(x: dict) -> str:
            try:
                return str((x or {}).get("path") or (x or {}).get("rel") or "").strip()
            except Exception:  # noqa: BLE001
                return ""

        def file_name(x: dict) -> str:
            try:
                return str((x or {}).get("name") or "").strip()
            except Exception:  # noqa: BLE001
                return ""

        def passes_ext(x: dict) -> bool:
            if not exts:
                return True
            n = file_name(x)
            p = file_path(x)
            cand = (Path(n).suffix or Path(p).suffix or "").lower()
            return bool(cand) and any(cand == e for e in exts)

        def split_alts(raw: str) -> List[str]:
            s = str(raw or "").strip()
            if not s:
                return []
            s = s.replace("（", "(").replace("）", ")")
            out = [x.strip() for x in s.split("/") if x.strip()]
            return out if out else [s]

        def derive_keyword_tokens(text: str) -> List[str]:
            t = str(text or "").strip()
            if not t:
                return []
            out: List[str] = []
            for m in re.finditer(r"[《“\"]([^》”\"]{2,50})[》”\"]", t):
                s = str(m.group(1) or "").strip()
                if s:
                    out.append(s)
            suffixes = "(?:说明书|安装包|报告|清单|方案|手册|纪要|记录表|意见|文档|材料|附件|交付物)"
            for m in re.finditer(rf"[\u4e00-\u9fff]{{2,25}}{suffixes}", t):
                out.append(str(m.group(0) or "").strip())
            for m in re.finditer(rf"(?:{suffixes})", t):
                out.append(str(m.group(0) or "").strip())
            for m in re.finditer(r"[A-Za-z]{2,}", t):
                out.append(str(m.group(0) or "").strip())
            seen = set()
            uniq: List[str] = []
            for x in out:
                k = norm(x)
                if not k or k in seen:
                    continue
                if len(k) < 2:
                    continue
                seen.add(k)
                uniq.append(x)
            return uniq[:12]

        def strict_missing_mode() -> bool:
            base = norm(checkpoint)
            if "《" in checkpoint and "》" in checkpoint:
                return True
            triggers = [
                "是否提供",
                "需提供",
                "应提供",
                "必须提供",
                "请提供",
                "提交",
                "交付",
                "附件",
                "缺失",
                "缺少",
                "未提供",
            ]
            for w in triggers:
                if norm(w) and norm(w) in base:
                    return True
            return False

        def allow_file_for_rule(file_text: str) -> bool:
            f = norm(file_text)
            if not f:
                return False
            r = norm(hint_text or checkpoint)
            review_file = any(k in f for k in [norm("联合评审意见"), norm("评审意见表"), norm("评审意见")])
            review_rule = any(k in r for k in [norm("评审"), norm("意见"), norm("纪要"), norm("会议"), norm("签字"), norm("确认")])
            if review_file and not review_rule:
                return False
            if review_rule:
                design_like = any(
                    k in f
                    for k in [
                        norm("概要设计说明书"),
                        norm("概要设计"),
                        norm("详细设计说明书"),
                        norm("详细设计"),
                        norm("系统需求规格说明书"),
                        norm("需求规格说明书"),
                        norm("数据库说明书"),
                        norm("数据库设计说明书"),
                    ]
                )
                if design_like:
                    return False
            ops_rule = any(k in r for k in [norm("运维"), norm("投产"), norm("上线"), norm("变更"), norm("发布")])
            if ops_rule:
                if any(
                    k in f
                    for k in [
                        norm("系统需求规格说明书"),
                        norm("需求规格说明书"),
                        norm("详细设计说明书"),
                        norm("详细设计"),
                        norm("数据库说明书"),
                        norm("数据库设计说明书"),
                        norm("概要设计说明书"),
                        norm("概要设计"),
                    ]
                ):
                    return False
            manual_file = any(k in f for k in [norm("用户操作手册"), norm("操作手册"), norm("用户手册"), norm("使用手册")])
            manual_rule = any(k in r for k in [norm("用户操作手册"), norm("操作手册"), norm("用户手册"), norm("使用手册")])
            if manual_file and not manual_rule:
                return False
            if forbidden_overlap(hint_text or checkpoint, file_text):
                return False
            return True

        def forbidden_overlap(rule_text: str, file_text: str) -> bool:
            r = norm(rule_text)
            f = norm(file_text)
            if not r or not f:
                return False
            pairs = [
                (["测试报告"], ["数据迁移方案", "迁移方案"]),
                (
                    ["测试", "测试计划", "测试用例", "用例", "testcase", "case", "缺陷", "bug", "回归"],
                    ["业务迁移方案", "业务差异分析报告", "业务差异解决方案", "业务需求规格说明书", "迁移方案", "差异分析", "差异解决", "需求规格说明书"],
                ),
                (["测试用例", "用例", "testcase", "case"], ["概要设计说明书", "概要设计", "详细设计说明书", "详细设计", "设计说明书"]),
                (["概要设计说明书", "概要设计"], ["评审意见", "联合评审意见表", "评审纪要", "联合评审意见"]),
                (["数据库说明书", "数据库设计", "数据库"], ["评审意见", "联合评审意见表", "评审纪要", "联合评审意见"]),
                (["系统需求规格说明书", "需求规格说明书", "srs"], ["评审意见", "联合评审意见表", "评审纪要", "联合评审意见"]),
            ]
            for rule_kws, file_kws in pairs:
                if any(norm(k) in r for k in rule_kws) and any(norm(k) in f for k in file_kws):
                    return True
            return False

        def tokens_from_pattern(pat: str) -> List[str]:
            s = str(pat or "").strip()
            if not s:
                return []
            parts = re.findall(r"[\u4e00-\u9fff]{2,}|[A-Za-z]{2,}|\d{2,}", s)
            stop = {
                "yyyy",
                "mm",
                "dd",
                "yyyymmdd",
                "vx",
                "xy",
                "vxy",
                "version",
                "ver",
            }
            out: List[str] = []
            for x in parts:
                k = norm(x)
                if not k:
                    continue
                if k in stop:
                    continue
                if re.fullmatch(r"\d{2,}", x):
                    continue
                out.append(x)
            return out

        def match_name_pattern(file_item: dict, name_pattern: str) -> bool:
            f_name = file_name(file_item)
            f_path = file_path(file_item)
            target = norm(f_path or f_name)
            if not target:
                return False
            for alt in split_alts(name_pattern):
                alt_suffix = Path(alt).suffix.lower()
                if alt_suffix:
                    if not (f_name.lower().endswith(alt_suffix) or f_path.lower().endswith(alt_suffix)):
                        continue
                toks = tokens_from_pattern(Path(alt).stem)
                if not toks:
                    continue
                ok = True
                for t in toks:
                    if norm(t) and norm(t) not in target:
                        ok = False
                        break
                if ok:
                    return True
            return False

        stage_filtered = [x for x in files if (not stage_need) or (file_stage(x) == stage_need)]
        stage_effective = stage_filtered if stage_filtered else files
        ext_filtered_raw = [x for x in stage_effective if passes_ext(x)]
        ext_filtered = [x for x in ext_filtered_raw if allow_file_for_rule(file_path(x) or file_name(x))]

        if exts and not ext_filtered_raw:
            return {
                "ok": True,
                "project": project,
                "stageNeed": stage_need,
                "extensions": exts,
                "source": "ext_missing",
                "counts": {"total": len(files), "afterStage": len(stage_effective), "afterExt": 0, "selected": 0},
                "candidates": [],
            }
        if ext_filtered_raw and not ext_filtered and strict_missing_mode():
            return {
                "ok": True,
                "project": project,
                "stageNeed": stage_need,
                "extensions": exts,
                "source": "restricted_missing",
                "counts": {"total": len(files), "afterStage": len(stage_effective), "afterExt": len(ext_filtered_raw), "selected": 0},
                "candidates": [],
            }

        selected: List[dict] = []
        selected_by = ""
        scope = str((rule or {}).get("scope") or "").strip()
        if naming:
            selected_by = "naming"
            for x in ext_filtered:
                for it in naming:
                    pat = str((it or {}).get("namePattern") or "").strip()
                    if not pat:
                        continue
                    if match_name_pattern(x, pat):
                        selected.append(x)
                        break
        if naming and not selected:
            if strict_missing_mode() or scope == "file":
                return {
                    "ok": True,
                    "project": project,
                    "stageNeed": stage_need,
                    "extensions": exts,
                    "source": "naming_missing",
                    "counts": {
                        "total": len(files),
                        "afterStage": len(stage_effective),
                        "afterExt": len(ext_filtered),
                        "selected": 0,
                    },
                    "candidates": [],
                }
            if ext_filtered:
                selected_by = "naming_fallback"
                selected = list(ext_filtered)

        if not selected:
            stop = {
                "是否",
                "检查",
                "抽查",
                "登记",
                "记录",
                "提供",
                "组织",
                "编写",
                "进行",
                "完成",
                "相关",
                "项目",
                "系统",
                "阶段",
                "过程",
                "材料",
                "文档",
                "文件",
                "附件",
                "内容",
                "要求",
                "情况",
                "工作",
                "结果",
                "说明",
                "需要",
                "必须",
                "应",
                "应当",
                "不得",
                "禁止",
                "可以",
            }

            raw_tokens: List[str] = []
            raw_tokens.extend([str(x or "").strip() for x in (keywords or []) if str(x or "").strip()])
            raw_tokens.extend(derive_keyword_tokens(checkpoint))
            for m in re.finditer(r"[\u4e00-\u9fff]{2,20}", checkpoint):
                raw_tokens.append(str(m.group(0) or "").strip())

            tokens: List[str] = []
            seen = set()
            for x in raw_tokens:
                k = norm(x)
                if not k:
                    continue
                if k in {norm(s) for s in stop}:
                    continue
                if len(k) < 3:
                    continue
                if k in seen:
                    continue
                seen.add(k)
                tokens.append(k)
            tokens = tokens[:12]

            best_score = 0
            scored: List[tuple[int, dict]] = []
            for f in ext_filtered:
                t = norm(file_path(f) or file_name(f))
                s = 0
                for k in tokens:
                    if k and k in t:
                        s += 1
                if s > 0:
                    scored.append((s, f))
                    if s > best_score:
                        best_score = s

            if best_score <= 0:
                selected_by = "token_missing"
                selected = []
            else:
                selected_by = "tokens"
                scored.sort(key=lambda x: x[0], reverse=True)
                selected = [f for s, f in scored if s == best_score]

        if not selected and not strict_missing_mode() and scope != "file":
            # 内容检查类规则：若无法从文件名/路径推导到目标交付物，则降级为对同阶段/后缀文件进行内容检查，
            # 避免“所有规则都缺少文件”导致无法对已有文件进行任何打开与检查。
            if ext_filtered:
                selected_by = "fallback"
                selected = list(ext_filtered)
            else:
                selected_by = "no_files"
                selected = []

        out = []
        for x in selected[:50]:
            out.append({"path": file_path(x), "name": file_name(x), "stage": file_stage(x)})

        return {
            "ok": True,
            "project": project,
            "stageNeed": stage_need,
            "extensions": exts,
            "source": selected_by or "unknown",
            "counts": {"total": len(files), "afterStage": len(stage_filtered), "afterExt": len(ext_filtered), "selected": len(selected)},
            "candidates": out,
        }

    @app.post("/api/scan/llm")
    async def scan_llm(request: Request):
        if not llm_configured():
            return JSONResponse(status_code=400, content={"ok": False, "message": "LLM 未配置：AUTOQA_LLM_API_KEY / OPENAI_API_KEY"})

        body = await request.json()
        project = str((body or {}).get("project") or "").strip()
        rule = (body or {}).get("rule") or {}
        candidates = (body or {}).get("candidates") or []
        if not isinstance(candidates, list):
            candidates = []

        rule_prompt = str(rule.get("prompt") or "").strip()
        scope = "文件名" if str(rule.get("scope") or "").strip() == "file" else "内容"
        stage = str(rule.get("stage") or "").strip() or "-"
        rule_name = str(rule.get("name") or "").strip()
        checkpoint = str(rule.get("checkpoint") or rule.get("pattern") or "").strip()
        if not rule_prompt:
            rule_prompt = f"【{stage}】\n检查要点（规则详细）：{checkpoint}\n检查对象：{scope}\n"

        file_list_lines: List[str] = []
        file_blocks: List[str] = []
        for x in candidates[:12]:
            p = str((x or {}).get("path") or "").strip()
            n = str((x or {}).get("name") or "").strip()
            file_list_lines.append(f"- {n or p}".strip())
            if scope == "内容":
                c = (x or {}).get("content")
                if c is None:
                    continue
                text = clamp_text(c, 12000)
                file_blocks.append(f"【{n or p}】\n{text}")

        extra = "\n\n".join(
            [
                f"项目：{project}" if project else "",
                f"规则：{rule_name}" if rule_name else "",
                "候选文件清单（仅用于定位，不要仅凭目录判断）：\n" + "\n".join(file_list_lines) if file_list_lines else "候选文件清单为空",
                "候选文件正文（可能截断）：\n\n" + "\n\n".join(file_blocks) if file_blocks else "",
            ]
        ).strip()

        system_msg = (
            "你是一个严谨的质量检查模型。请严格依据“检查要点/检查对象”对候选文件进行判断。"
            "只输出一个 JSON 对象，不要输出任何额外文字。"
            'JSON Schema: {"conclusion":"符合|不符合","reason":"...","suggestion":"...","evidence":[{"file":"...","snippet":"..."}]}'
        )
        user_msg = (rule_prompt + "\n\n" + extra).strip()

        lg = logging.getLogger("autoqa.scan")
        lg.info("scan_llm start project=%s rule=%s stage=%s scope=%s candidates=%s", project or "-", rule_name or "-", stage or "-", scope, len(candidates))
        try:
            content = call_llm([{"role": "system", "content": system_msg}, {"role": "user", "content": user_msg}])
        except Exception as e:  # noqa: BLE001
            lg.exception("scan_llm error project=%s rule=%s: %s", project or "-", rule_name or "-", str(e))
            raise

        parsed = parse_json_maybe(content)
        if not parsed:
            parsed = {
                "conclusion": "不符合",
                "reason": "模型输出未按 JSON 格式返回，无法解析",
                "suggestion": "",
                "evidence": [{"file": "", "snippet": clamp_text(content, 800)}],
            }
        lg.info("scan_llm done project=%s rule=%s conclusion=%s evidence=%s", project or "-", rule_name or "-", str(parsed.get("conclusion") or ""), len(parsed.get("evidence") or []))
        return {"ok": True, "result": parsed, "raw": clamp_text(content, 4000)}

    @app.post("/api/extract/text")
    async def extract_text_api(request: Request):
        import tempfile

        lg = logging.getLogger("autoqa.extract")
        filename = request.headers.get("x-filename") or request.query_params.get("filename") or "file"
        suffix = Path(str(filename)).suffix.lower()
        data = await request.body()
        lg.info("extract_text start filename=%s size=%s", str(filename), len(data or b""))

        tmp_dir = (work_inputs_dir / "_tmp").resolve()
        tmp_dir.mkdir(parents=True, exist_ok=True)
        tmp_path = None
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix, dir=str(tmp_dir)) as fp:
                fp.write(data or b"")
                tmp_path = Path(fp.name)
            text = extract_text(tmp_path)
            if text is None or not str(text).strip():
                lg.info("extract_text unreadable filename=%s", str(filename))
                return {"ok": True, "text": "", "unreadable": True}
            lg.info("extract_text ok filename=%s chars=%s", str(filename), len(text))
            return {"ok": True, "text": str(text), "unreadable": False}
        except Exception as e:  # noqa: BLE001
            lg.exception("extract_text error filename=%s: %s", str(filename), str(e))
            return JSONResponse(status_code=500, content={"ok": False, "message": str(e), "filename": str(filename)})
        finally:
            try:
                if tmp_path and tmp_path.exists():
                    tmp_path.unlink()
            except Exception:  # noqa: BLE001
                pass

    @app.post("/api/rules/upload")
    async def upload_rules(request: Request):
        filename = request.headers.get("x-filename") or request.query_params.get("filename") or "rules.xlsx"
        suffix = Path(filename).suffix.lower()
        if suffix not in {".xlsx", ".xls"}:
            return JSONResponse(status_code=400, content={"ok": False, "message": "仅支持 .xls/.xlsx"})

        dst = work_inputs_dir / "rules.xlsx"
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        backup = work_inputs_dir / f"rules__{stamp}{suffix or '.xlsx'}"

        data = await request.body()
        dst.write_bytes(data)
        try:
            backup.write_bytes(data)
        except Exception:  # noqa: BLE001
            pass

        return {
            "ok": True,
            "path": str(dst),
            "backup": str(backup),
        }

    @app.post("/api/rules/parsed")
    async def save_parsed_rules(payload: List[dict]):
        out_path = work_rules_dir / "rules.json"
        out_path.write_text(json.dumps(payload or [], ensure_ascii=False, indent=2), encoding="utf-8")
        return {"ok": True, "path": str(out_path), "count": len(payload or [])}

    @app.get("/api/rules/parsed")
    async def load_parsed_rules():
        out_path = work_rules_dir / "rules.json"
        if not out_path.exists() or not out_path.is_file():
            return {"ok": True, "path": str(out_path), "rules": [], "count": 0, "exists": False}
        try:
            payload = json.loads(out_path.read_text(encoding="utf-8"))
        except Exception as e:  # noqa: BLE001
            return JSONResponse(status_code=500, content={"ok": False, "message": str(e), "path": str(out_path)})
        if not isinstance(payload, list):
            payload = []
        return {"ok": True, "path": str(out_path), "rules": payload, "count": len(payload), "exists": True}

    @app.delete("/api/rules/parsed")
    async def delete_parsed_rules():
        out_path = work_rules_dir / "rules.json"
        deleted = False
        try:
            if out_path.exists():
                out_path.unlink()
                deleted = True
        except Exception as e:  # noqa: BLE001
            return JSONResponse(status_code=500, content={"ok": False, "message": str(e), "path": str(out_path)})
        return {"ok": True, "path": str(out_path), "deleted": deleted}

    @app.post("/api/config/upload")
    async def upload_config(request: Request):
        filename = request.headers.get("x-filename") or request.query_params.get("filename") or "config.xlsx"
        suffix = Path(filename).suffix.lower()
        if suffix not in {".xlsx", ".xls"}:
            return JSONResponse(status_code=400, content={"ok": False, "message": "仅支持 .xls/.xlsx"})

        dst = work_inputs_dir / "config.xlsx"
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        backup = work_inputs_dir / f"config__{stamp}{suffix or '.xlsx'}"

        data = await request.body()
        dst.write_bytes(data)
        try:
            backup.write_bytes(data)
        except Exception:  # noqa: BLE001
            pass

        return {
            "ok": True,
            "path": str(dst),
            "backup": str(backup),
        }

    @app.post("/api/prompts/generate")
    async def generate_prompts(payload: List[dict]):
        rules_in = payload or []
        rules_out = [generate_prompt_from_skill(r) for r in rules_in]
        return {"ok": True, "rules": rules_out, "count": len(rules_out), "skill": str(skill_prompt_path)}

    return app



def main(argv: Optional[Sequence[str]] = None) -> int:
    repo_root = Path(__file__).resolve().parents[1]
    default_input_dir = repo_root / "work" / "scan"
    default_rules_excel = repo_root / "work" / "inputs" / "rules.xlsx"
    default_report_xlsx = repo_root / "work" / "outputs" / "report.xlsx"
    default_rules_json = repo_root / "work" / "rules" / "rules.json"
    default_config_yaml = repo_root / "scripts" / "config.yaml"

    parser = argparse.ArgumentParser()
    parser.add_argument("--serve", action="store_true", help="start API server for UI persistence")
    parser.add_argument("--host", default="0.0.0.0")
    parser.add_argument("--port", type=int, default=8000)
    parser.add_argument("--config", default=str(default_config_yaml), help="path to config.yaml")
    parser.add_argument("--input", default=str(default_input_dir), help="root directory to scan")
    parser.add_argument("--rules", default=str(default_rules_excel), help="rules excel (.xlsx)")
    parser.add_argument("--output", default=str(default_report_xlsx), help="output report excel (.xlsx)")
    parser.add_argument("--parsed-rules-json", default=str(default_rules_json), help="write parsed rules json")
    parser.add_argument("--no-save-parsed-rules", action="store_true", help="do not write parsed rules to work/rules")
    parser.add_argument("--no-persist-inputs", action="store_true", help="do not persist input artifacts to work/inputs")
    parser.add_argument("--max-mb", type=int, default=20)
    parser.add_argument(
        "--ext",
        action="append",
        default=None,
        help="restrict scanned file extensions, e.g. --ext .js --ext .pdf",
    )

    args = parser.parse_args(argv)

    if args.serve:
        import uvicorn
        import logging

        cfg_path = Path(str(args.config or "")).expanduser().resolve()
        app = create_api_app(repo_root, config_path=cfg_path)
        log_dir = (repo_root / "work" / "logs").resolve()
        log_dir.mkdir(parents=True, exist_ok=True)
        log_path = log_dir / f"backend_{datetime.now().strftime('%Y%m%d_%H%M%S')}.log"

        root_logger = logging.getLogger()
        root_logger.setLevel(logging.INFO)
        for h in list(root_logger.handlers):
            try:
                root_logger.removeHandler(h)
            except Exception:  # noqa: BLE001
                pass

        fmt = logging.Formatter("%(asctime)s %(levelname)s %(name)s: %(message)s")
        file_handler = logging.FileHandler(str(log_path), encoding="utf-8")
        file_handler.setLevel(logging.INFO)
        file_handler.setFormatter(fmt)
        stream_handler = logging.StreamHandler(sys.stdout)
        stream_handler.setLevel(logging.INFO)
        stream_handler.setFormatter(fmt)
        root_logger.addHandler(file_handler)
        root_logger.addHandler(stream_handler)

        for name in ("uvicorn", "uvicorn.error", "uvicorn.access"):
            lg = logging.getLogger(name)
            lg.handlers = []
            lg.propagate = True

        sys.stdout.write(f"Log file: {log_path}\n")
        uvicorn.run(app, host=args.host, port=args.port, log_level="info", access_log=True, log_config=None)
        return 0

    root = Path(args.input).expanduser().resolve()
    rules_path = Path(args.rules).expanduser().resolve()
    out_path = Path(args.output).expanduser().resolve()

    work_rules_dir = (repo_root / "work" / "rules").resolve()
    work_outputs_dir = (repo_root / "work" / "outputs").resolve()
    work_scan_dir = (repo_root / "work" / "scan").resolve()
    work_inputs_dir = (repo_root / "work" / "inputs").resolve()

    work_rules_dir.mkdir(parents=True, exist_ok=True)
    work_outputs_dir.mkdir(parents=True, exist_ok=True)
    work_scan_dir.mkdir(parents=True, exist_ok=True)
    work_inputs_dir.mkdir(parents=True, exist_ok=True)
    if not root.exists() or not root.is_dir():
        raise RuntimeError(f"Input dir not found: {root}")
    if not rules_path.exists() or not rules_path.is_file():
        raise RuntimeError(f"Rules file not found: {rules_path}")

    rules = load_rules_from_excel(rules_path)
    findings = scan(root, rules, include_exts=args.ext, max_mb=args.max_mb)

    if not args.no_persist_inputs:
        dst = work_inputs_dir / "rules.xlsx"
        try:
            if rules_path.resolve() != dst.resolve():
                shutil.copy2(rules_path, dst)
        except Exception:  # noqa: BLE001
            pass

        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        backup = work_inputs_dir / f"rules__{stamp}{rules_path.suffix.lower() or '.xlsx'}"
        try:
            shutil.copy2(rules_path, backup)
        except Exception:  # noqa: BLE001
            pass

    if not args.no_save_parsed_rules:
        rules_json_path.parent.mkdir(parents=True, exist_ok=True)
        write_rules_json(rules_json_path, rules)
    write_report_xlsx(out_path, findings)
    sys.stdout.write(f"Wrote: {out_path}\n")
    sys.stdout.write(f"Findings: {len(findings)}\n")
    return 0
    if not args.no_save_parsed_rules:
        sys.stdout.write(f"Parsed rules: {rules_json_path}\n")




if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as e:  # noqa: BLE001
        sys.stderr.write(str(e) + "\n")
        raise SystemExit(1)
